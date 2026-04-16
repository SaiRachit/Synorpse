import os
import json
import tempfile
from pathlib import Path
from groq import Groq
from dotenv import load_dotenv
import PyPDF2
from docx import Document
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import io
import hashlib
from typing import List, Dict
import pickle

load_dotenv()

_PROJECT_ROOT = Path(__file__).resolve().parent.parent


def _find_poppler_path():
    """Return a Poppler bin directory if one is available."""
    env_path = os.environ.get("POPPLER_PATH")
    if env_path and Path(env_path).exists():
        return env_path

    bundled_bin = _PROJECT_ROOT / "tools" / "poppler"
    matches = sorted(bundled_bin.rglob("pdftoppm.exe"))
    if matches:
        return str(matches[0].parent)

    return None


POPPLER_PATH = _find_poppler_path()

class DocumentQA:
    def __init__(self, api_key=None, chunk_size=3000, chunk_overlap=200):
        """
        Initialize with Groq API key from .env file
        
        Args:
            api_key: Groq API key (optional, will use .env if not provided)
            chunk_size: Size of text chunks in characters
            chunk_overlap: Overlap between chunks to maintain context
        """
        self.api_key = api_key or os.environ.get("GroqAPIKeyDoc")
        if not self.api_key:
            raise ValueError("GroqAPIKeyDoc not found in .env file. Please create a .env file with GroqAPIKeyDoc=your-key-here")
        
        self.client = Groq(api_key=self.api_key)
        self.metadata_file = None
        self.document_content = None
        self.chunks = []
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.document_hash = None
        self.cache_dir = Path(tempfile.gettempdir()) / "doc_qa_cache"
        self.cache_dir.mkdir(exist_ok=True)
    
    def read_pdf_with_ocr(self, file_path, start_page=None, end_page=None):
        """
        Read PDF with OCR for scanned/image-based PDFs
        Memory efficient - processes pages in batches
        """
        try:
            print("Attempting OCR on PDF (this may take a moment)...")
            # Get total pages first
            images = convert_from_path(
                file_path,
                first_page=start_page,
                last_page=end_page,
                poppler_path=POPPLER_PATH
            )
            text_content = []
            
            # Process in smaller batches to save memory
            batch_size = 5
            for i in range(0, len(images), batch_size):
                batch = images[i:i+batch_size]
                for j, image in enumerate(batch):
                    page_num = i + j + (start_page or 1)
                    print(f"  Processing page {page_num}...")
                    text = pytesseract.image_to_string(image)
                    text_content.append(text)
                    # Clear image from memory
                    image.close()
                # Clear batch from memory
                del batch
            
            return '\n\n'.join(text_content)
        except Exception as e:
            print(f"  OCR failed: {e}")
            return None
    
    def read_pdf_streaming(self, file_path):
        """Read PDF file in streaming fashion for large files"""
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                total_pages = len(pdf_reader.pages)
                print(f"  Total pages: {total_pages}")
                
                # For very large PDFs, process in batches
                batch_size = 50
                all_text = []
                
                for i in range(0, total_pages, batch_size):
                    end = min(i + batch_size, total_pages)
                    print(f"  Reading pages {i+1}-{end}...")
                    
                    batch_text = []
                    for page_num in range(i, end):
                        page = pdf_reader.pages[page_num]
                        text = page.extract_text()
                        batch_text.append(text)
                    
                    all_text.extend(batch_text)
                
                extracted_text = '\n\n'.join(all_text)
                
                # Check if OCR is needed
                if len(extracted_text.strip()) < 100:
                    print("  Text extraction yielded minimal content, trying OCR...")
                    ocr_text = self.read_pdf_with_ocr(file_path)
                    if ocr_text:
                        return ocr_text, total_pages
                
                return extracted_text, total_pages
        except Exception as e:
            print(f"  Standard PDF reading failed: {e}")
            print("  Attempting OCR fallback...")
            ocr_text = self.read_pdf_with_ocr(file_path)
            if ocr_text:
                pages = len(convert_from_path(file_path, poppler_path=POPPLER_PATH))
                return ocr_text, pages
            raise
    
    def read_word(self, file_path):
        """Read Word document (.docx) with memory efficiency"""
        doc = Document(file_path)
        text_content = []
        
        # Process paragraphs in batches
        batch_size = 100
        for i in range(0, len(doc.paragraphs), batch_size):
            batch = doc.paragraphs[i:i+batch_size]
            text_content.extend([p.text for p in batch if p.text.strip()])
            if i % 500 == 0 and i > 0:
                print(f"  Processed {i} paragraphs...")
        
        # Process tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_content.append(row_text)
        
        return '\n'.join(text_content)
    
    def read_powerpoint(self, file_path):
        """Read PowerPoint presentation (.pptx)"""
        prs = Presentation(file_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides):
            if i % 50 == 0 and i > 0:
                print(f"  Processed {i} slides...")
            text_content.append(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
        
        return '\n'.join(text_content)
    
    def read_image(self, file_path):
        """Read text from image using OCR"""
        image = Image.open(file_path)
        # Resize very large images to save memory
        max_size = 4000
        if image.width > max_size or image.height > max_size:
            print(f"  Resizing large image for OCR...")
            image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        text = pytesseract.image_to_string(image)
        image.close()
        return text
    
    def compute_file_hash(self, file_path):
        """Compute hash of file for caching"""
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            # Read in chunks to handle large files
            for chunk in iter(lambda: f.read(8192), b''):
                hasher.update(chunk)
        return hasher.hexdigest()
    
    def load_from_cache(self, file_path):
        """Load preprocessed document from cache if available"""
        file_hash = self.compute_file_hash(file_path)
        cache_file = self.cache_dir / f"{file_hash}.pkl"
        
        if cache_file.exists():
            try:
                print(f"Loading from cache...")
                with open(cache_file, 'rb') as f:
                    cached_data = pickle.load(f)
                self.document_content = cached_data['content']
                self.chunks = cached_data['chunks']
                self.document_hash = file_hash
                print(f" Loaded from cache: {len(self.chunks)} chunks")
                return True
            except Exception as e:
                print(f"Cache load failed: {e}")
                return False
        return False
    
    def save_to_cache(self, file_path):
        """Save preprocessed document to cache"""
        if not self.document_content:
            return False
        
        try:
            file_hash = self.compute_file_hash(file_path)
            cache_file = self.cache_dir / f"{file_hash}.pkl"
            
            cache_data = {
                'content': self.document_content,
                'chunks': self.chunks,
                'file_path': file_path
            }
            
            with open(cache_file, 'wb') as f:
                pickle.dump(cache_data, f)
            
            self.document_hash = file_hash
            print(f" Saved to cache for faster future access")
            return True
        except Exception as e:
            print(f"Cache save failed: {e}")
            return False
    
    def chunk_text(self, text: str) -> List[Dict]:
        """
        Split text into overlapping chunks for better context handling
        Returns list of dicts with chunk text and metadata
        """
        chunks = []
        words = text.split()
        
        # Calculate approximate chunk size in words
        avg_word_length = sum(len(word) for word in words[:1000]) / min(1000, len(words))
        chunk_words = int(self.chunk_size / avg_word_length)
        overlap_words = int(self.chunk_overlap / avg_word_length)
        
        print(f"Creating chunks (size: {self.chunk_size} chars, overlap: {self.chunk_overlap} chars)...")
        
        for i in range(0, len(words), chunk_words - overlap_words):
            chunk_words_list = words[i:i + chunk_words]
            chunk_text = ' '.join(chunk_words_list)
            
            chunks.append({
                'text': chunk_text,
                'start_word': i,
                'end_word': i + len(chunk_words_list),
                'chunk_id': len(chunks)
            })
            
            if len(chunks) % 100 == 0:
                print(f"  Created {len(chunks)} chunks...")
        
        print(f" Created {len(chunks)} chunks")
        return chunks
    
    def read(self, file_path):
        """Compatibility method for DocumentReader"""
        return self.read_document(file_path)

    def read_document(self, file_path, use_cache=True):
        """Read the entire document and store content (supports multiple formats)"""
        # Try loading from cache first
        if use_cache and self.load_from_cache(file_path):
            return True
        
        try:
            file_extension = Path(file_path).suffix.lower()
            file_size = os.path.getsize(file_path) / (1024 * 1024)  # Size in MB
            print(f"Reading {file_extension} file ({file_size:.2f} MB)...")
            
            if file_extension == '.pdf':
                content, pages = self.read_pdf_streaming(file_path)
                self.document_content = content
                print(f" PDF loaded: {pages} pages, {len(self.document_content):,} characters")
                
            elif file_extension in ['.docx', '.doc']:
                self.document_content = self.read_word(file_path)
                print(f" Word document loaded: {len(self.document_content):,} characters")
                
            elif file_extension in ['.pptx', '.ppt']:
                self.document_content = self.read_powerpoint(file_path)
                print(f" PowerPoint loaded: {len(self.document_content):,} characters")
                
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
                self.document_content = self.read_image(file_path)
                print(f" Image text extracted: {len(self.document_content):,} characters")
                
            elif file_extension in ['.txt', '.md', '.csv', '.log']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    self.document_content = f.read()
                print(f" Text file loaded: {len(self.document_content):,} characters")
                
            else:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        self.document_content = f.read()
                    print(f" File loaded as text: {len(self.document_content):,} characters")
                except:
                    raise ValueError(f"Unsupported file format: {file_extension}")
            
            if not self.document_content or len(self.document_content.strip()) < 10:
                print(" Warning: Very little content was extracted from the document")
                return False
            
            # Create chunks for large documents
            if len(self.document_content) > 10000:
                self.chunks = self.chunk_text(self.document_content)
            else:
                # For small documents, use single chunk
                self.chunks = [{
                    'text': self.document_content,
                    'start_word': 0,
                    'end_word': len(self.document_content.split()),
                    'chunk_id': 0
                }]
            
            # Save to cache
            if use_cache:
                self.save_to_cache(file_path)
            
            return True
            
        except Exception as e:
            print(f" Error reading document: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def find_relevant_chunks(self, question: str, top_k: int = 3) -> List[Dict]:
        """
        Find most relevant chunks for the question using simple keyword matching
        For better results, could use embeddings (would require additional API)
        """
        if len(self.chunks) <= 3:
            return self.chunks
        
        # Simple keyword-based relevance scoring
        question_words = set(question.lower().split())
        scored_chunks = []
        
        for chunk in self.chunks:
            chunk_words = set(chunk['text'].lower().split())
            # Calculate overlap
            overlap = len(question_words & chunk_words)
            scored_chunks.append((chunk, overlap))
        
        # Sort by score and return top_k
        scored_chunks.sort(key=lambda x: x[1], reverse=True)
        return [chunk for chunk, score in scored_chunks[:top_k]]
    
    def save_metadata(self):
        """Save document metadata temporarily"""
        if not self.document_content:
            print(" No document content to save")
            return False
        
        try:
            temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8')
            self.metadata_file = temp_file.name
            
            metadata = {
                'length': len(self.document_content),
                'word_count': len(self.document_content.split()),
                'num_chunks': len(self.chunks),
                'chunk_size': self.chunk_size
            }
            
            json.dump(metadata, temp_file)
            temp_file.close()
            
            print(f" Metadata saved: {metadata['word_count']:,} words in {metadata['num_chunks']} chunks")
            return True
        except Exception as e:
            print(f" Error saving metadata: {e}")
            return False
    
    def ask_question(self, question, use_all_context=False):
        """
        Ask a question about the document using Groq API
        
        Args:
            question: The question to ask
            use_all_context: If True, tries to use entire document (may fail for very large docs)
        """
        if not self.document_content:
            return "No document loaded. Please load a document first."
        
        try:
            # For small documents or if requested, use full context
            if use_all_context or len(self.chunks) <= 1:
                max_chars = 100000  # Increased limit
                content = self.document_content[:max_chars]
                if len(self.document_content) > max_chars:
                    content += "\n\n[Document truncated due to length...]"
                context_source = "full document (truncated)" if len(self.document_content) > max_chars else "full document"
            else:
                # Use relevant chunks for large documents
                relevant_chunks = self.find_relevant_chunks(question, top_k=5)
                content = "\n\n---\n\n".join([chunk['text'] for chunk in relevant_chunks])
                context_source = f"{len(relevant_chunks)} relevant sections"
            
            print(f"Using context from: {context_source}")
            
            messages = [
                {
                    "role": "system",
                    "content": """You are a helpful assistant that answers questions about documents. 
                    Provide accurate, detailed answers based on the document content provided.
                    If the answer is not in the provided context, say so clearly.
                    Always cite specific parts of the document when possible."""
                },
                {
                    "role": "user",
                    "content": f"Document content:\n\n{content}\n\nQuestion: {question}"
                }
            ]
            
            response = self.client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages,
                temperature=0.3,
                max_tokens=2048  # Increased for more detailed answers
            )
            
            answer = response.choices[0].message.content
            return answer
            
        except Exception as e:
            return f"Error getting answer: {e}"
    
    def get_document_summary(self):
        """Generate a summary of the entire document"""
        if not self.document_content:
            return "No document loaded."
        
        try:
            # For very large documents, summarize in parts then combine
            if len(self.chunks) > 10:
                print("Generating summary from document chunks...")
                chunk_summaries = []
                
                # Summarize every Nth chunk to get representative sample
                step = max(1, len(self.chunks) // 10)
                for i in range(0, len(self.chunks), step):
                    chunk = self.chunks[i]
                    print(f"  Summarizing chunk {i+1}/{len(self.chunks)}...")
                    
                    messages = [
                        {
                            "role": "system",
                            "content": "Provide a brief summary of the key points in this text section."
                        },
                        {
                            "role": "user",
                            "content": chunk['text'][:5000]  # Limit chunk size
                        }
                    ]
                    
                    response = self.client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=messages,
                        temperature=0.3,
                        max_tokens=512
                    )
                    
                    chunk_summaries.append(response.choices[0].message.content)
                
                # Combine chunk summaries
                combined = "\n\n".join(chunk_summaries)
                
                messages = [
                    {
                        "role": "system",
                        "content": "Create a comprehensive summary from these section summaries."
                    },
                    {
                        "role": "user",
                        "content": f"Section summaries:\n\n{combined}"
                    }
                ]
                
            else:
                # For smaller documents, summarize directly
                content = self.document_content[:50000]
                messages = [
                    {
                        "role": "system",
                        "content": "Provide a comprehensive and structured summary of this document. "
                                   "Highlight key points, main topics, and important conclusions. "
                                   "If the document contains instructions, laboratories, or guides, "
                                   "explicitly list the distinct Steps, Tasks, or Laboratories mentioned."
                    },
                    {
                        "role": "user",
                        "content": content
                    }
                ]
            
            response = self.client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages,
                temperature=0.3,
                max_tokens=2048
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            return f"Error generating summary: {e}"
    
    def cleanup(self):
        """Delete temporary metadata file (keep cache for reuse)"""
        if self.metadata_file and os.path.exists(self.metadata_file):
            try:
                os.remove(self.metadata_file)
                print(f" Temporary metadata cleaned up")
                self.metadata_file = None
                self.document_content = None
                self.chunks = []
                return True
            except Exception as e:
                print(f" Error cleaning up: {e}")
                return False
        else:
            self.document_content = None
            self.chunks = []
            return True
    
    def clear_cache(self):
        """Clear all cached documents"""
        try:
            import shutil
            if self.cache_dir.exists():
                shutil.rmtree(self.cache_dir)
                self.cache_dir.mkdir(exist_ok=True)
                print(" Cache cleared")
            return True
        except Exception as e:
            print(f" Error clearing cache: {e}")
            return False

def main():
    """Main function to run the document Q&A system"""
    print("=== Improved Universal Document Q&A System with Groq API ===")
    print(" Supports very large documents with intelligent chunking")
    print(" Caching for faster repeat access")
    print(" Memory-efficient processing")
    print("Supports: PDF, Word, PowerPoint, Images (with OCR), and text files\n")
    
    # Initialize
    try:
        qa_system = DocumentQA(chunk_size=3000, chunk_overlap=200)
    except ValueError as e:
        print(f"Error: {e}")
        print("\nTo use this script:")
        print("1. Get a Groq API key from https://console.groq.com")
        print("2. Create a .env file in the same directory as this script")
        print("3. Add this line to the .env file: GroqAPIKeyDoc=your-key-here")
        print("\n4. Install required packages:")
        print("   pip install groq python-dotenv PyPDF2 python-docx python-pptx pytesseract Pillow pdf2image")
        print("\n5. Install Tesseract OCR:")
        print("   - Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki")
        print("   - Mac: brew install tesseract")
        print("   - Linux: sudo apt-get install tesseract-ocr")
        print("\n6. Install poppler (for PDF to image conversion):")
        print("   - Windows: Download from https://github.com/oschwartz10612/poppler-windows/releases")
        print("   - Mac: brew install poppler")
        print("   - Linux: sudo apt-get install poppler-utils")
        return
    
    doc_path = input("Enter the path to your document: ").strip().strip('"')
    
    if not os.path.exists(doc_path):
        print(f" File not found: {doc_path}")
        return
    
    print()
    if not qa_system.read_document(doc_path):
        return
    
    if not qa_system.save_metadata():
        return
    
    print("\n--- You can now ask questions about the document ---")
    print("Commands:")
    print("  - Type your question to ask")
    print("  - 'summary' - Get document summary")
    print("  - 'stats' - Show document statistics")
    print("  - 'quit' or 'exit' - Finish and cleanup")
    print()
    
    # Q&A loop
    while True:
        question = input("\nYour question: ").strip()
        
        if question.lower() in ['quit', 'exit', 'q']:
            break
        
        if not question:
            continue
        
        if question.lower() == 'summary':
            print("\nGenerating document summary...")
            summary = qa_system.get_document_summary()
            print(f"\nSummary:\n{summary}")
            continue
        
        if question.lower() == 'stats':
            print(f"\nDocument Statistics:")
            print(f"  Total characters: {len(qa_system.document_content):,}")
            print(f"  Total words: {len(qa_system.document_content.split()):,}")
            print(f"  Number of chunks: {len(qa_system.chunks)}")
            print(f"  Chunk size: {qa_system.chunk_size} characters")
            continue
        
        print("\nThinking...")
        answer = qa_system.ask_question(question)
        print(f"\nAnswer: {answer}")
    
    print("\n--- Cleaning up ---")
    qa_system.cleanup()
    print("\n Session ended. Temporary data deleted (cache preserved for future use).")
    print(f"To clear cache, delete: {qa_system.cache_dir}")

# Alias for backward compatibility and standard naming
DocumentReader = DocumentQA

if __name__ == "__main__":
    main()
